import os
import subprocess
import pandas as pd

from io import BytesIO
from pypdf import PdfReader, PdfWriter
from PIL import Image, ImageDraw, ImageFont
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from xml.etree import ElementTree as ET
from pptx import Presentation
from spire.doc import Document, TextWatermark, WatermarkLayout, Color

from app.celery_worker import celery_app


def create_watermark_pdf(watermark_text, page_width, page_height, font_size=50, opacity=0.3, rotation_angle=45):
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=(page_width, page_height))

    c.setFont("Helvetica-Bold", font_size)
    c.setFillGray(0, opacity)  # Black with transparency

    # Move to center, rotate, then draw text
    c.saveState()
    c.translate(page_width / 2, page_height / 2)
    c.rotate(rotation_angle)
    c.drawCentredString(0, 0, watermark_text)
    c.restoreState()

    c.save()
    buffer.seek(0)
    return buffer

def csv_to_pdf(csv_file, pdf_file):
    df = pd.read_csv(csv_file)

    c = canvas.Canvas(pdf_file, pagesize=A4)
    _, height = A4

    text = c.beginText(40, height - 40)
    text.setFont("Helvetica", 10)

    # Write header
    text.textLine(",".join(df.columns))

    # Write rows
    for _, row in df.iterrows():
        text.textLine(",".join(map(str, row.values)))

    c.drawText(text)
    c.save()

def create_transparent_watermark(text, font_size=40, opacity=80, angle=45, image_size=(200, 200)):
    
    # Transparent image (RGBA)
    img = Image.new("RGBA", image_size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)

    # Load font (better than load_default, which ignores size)
    font = ImageFont.load_default(size=font_size)

    # Measure text size
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width, text_height = bbox[2] - bbox[0], bbox[3] - bbox[1]

    # Center position
    x = (image_size[0] - text_width) // 2
    y = (image_size[1] - text_height) // 2

    # Transparent text (gray with alpha)
    text_color = (100, 100, 100, opacity)

    # Draw text
    draw.text((x, y), text, font=font, fill=text_color)

    # Rotate text
    img = img.rotate(angle, expand=1)

    # Save as PNG with transparency
    img.save("output/watermark.png", "PNG")

    return "output/watermark.png"

def convert_ppt_to_pptx(input_ppt):
    """Convert .ppt to .pptx using LibreOffice (headless)"""
    output_dir = os.path.dirname(os.path.abspath(input_ppt)) or "."
    try:
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pptx",
            "--outdir", output_dir, input_ppt
        ], check=True)
    except Exception as e:
        raise RuntimeError(f"Failed to convert {input_ppt} to pptx: {e}")

    base = os.path.splitext(input_ppt)[0]
    return base + ".pptx"

def add_watermark(input_pptx, output_pptx, watermark_text="DRAFT"):
    """Add watermark to pptx file"""
    prs = Presentation(input_pptx)
    watermark_file = create_transparent_watermark(watermark_text)

    for slide in prs.slides:
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        pic = slide.shapes.add_picture(watermark_file, 0, 0)
        pic.left = int((slide_width - pic.width) / 2)
        pic.top = int((slide_height - pic.height) / 2)

    prs.save(output_pptx)

@celery_app.task
def add_watermark_to_pdf(input_pdf: str, output_pdf: str, watermark_text='DRAFT'):

    input_pdf_reader = PdfReader(open(input_pdf, "rb"))
    output_pdf_writer = PdfWriter()

    for page in input_pdf_reader.pages:
        page_width = float(page.mediabox.width)
        page_height = float(page.mediabox.height)

        # Create watermark matching this page size
        watermark_pdf_buffer = create_watermark_pdf(watermark_text, page_width, page_height)
        watermark_reader = PdfReader(watermark_pdf_buffer)
        watermark_page = watermark_reader.pages[0]

        # Merge watermark into current page
        page.merge_page(watermark_page)
        output_pdf_writer.add_page(page)

    with open(output_pdf, "wb") as output_file:
        output_pdf_writer.write(output_file)

@celery_app.task
def add_watermark_to_image(image_path, output_path, watermark_text="DRAFT"):

    img = Image.open(image_path).convert("RGBA")
    img_width, img_height = img.size

    font_size_ratio=0.1
    color=(150, 150, 150, 100)
    font_size = int(img_height * font_size_ratio)
    font = ImageFont.load_default(size=font_size)

    # Transparent layer for text
    text_layer = Image.new("RGBA", (img_width, img_height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(text_layer)

    # Get text size
    text_bbox = draw.textbbox((0, 0), watermark_text, font=font)
    text_width = text_bbox[2] - text_bbox[0]
    text_height = text_bbox[3] - text_bbox[1]

    # Center position
    x = (img_width - text_width) / 2
    y = (img_height - text_height) / 2

    draw.text((x, y), watermark_text, font=font, fill=color)

    # Rotate
    rotated_text = text_layer.rotate(45, expand=False)

    # Overlay
    watermarked = Image.alpha_composite(img, rotated_text)

    if output_path.lower().endswith((".jpg", ".jpeg")):
        watermarked = watermarked.convert("RGB")
    
    watermarked.save(output_path)

@celery_app.task
def add_watermark_to_rtf(rtf_path, output_path, watermark_text="DRAFT"):
        
    doc = Document()
    doc.LoadFromFile(rtf_path)

    # Create watermark text
    watermark = TextWatermark()
    watermark.Text = "DRAFT"
    watermark.FontSize = 60                # adjust size as needed
    watermark.Color = Color.get_LightGray()
    watermark.Layout = WatermarkLayout.Diagonal  # diagonal alignment
    watermark.FontName = "Arial"

    # Apply watermark
    doc.Watermark = watermark

    # Save to RTF
    doc.SaveToFile(output_path)

@celery_app.task
def add_watermark_to_docx(input_doc_path, output_doc_path, watermark_text="DRAFT"):

    # Load the document
    document = Document()
    document.LoadFromFile(input_doc_path)

    # Create a text watermark
    watermark = TextWatermark()
    watermark.Text = watermark_text
    watermark.FontName = "Arial"
    watermark.FontSize = 72
    watermark.Color = Color.get_LightGray()
    watermark.Layout = WatermarkLayout.Diagonal # Set diagonal layout

    # If the document has at least one paragraph, delete the first one
    if document.Sections.Count > 0:
        first_section = document.Sections[0]
        first_paragraph = first_section.Paragraphs[0]
        first_section.Paragraphs.Remove(first_paragraph)

    # Apply the watermark
    document.Watermark = watermark

    # Save the document
    document.SaveToFile(output_doc_path)

@celery_app.task
def add_watermark_to_csv(input_csv_path, output_csv_to_pdf, watermark_text="DRAFT"):
    
    csv_to_pdf_path = "output/csv_to_pdf.pdf"
    csv_to_pdf(input_csv_path, csv_to_pdf_path)

    add_watermark_to_pdf(csv_to_pdf_path, output_csv_to_pdf)

@celery_app.task
def add_watermark_to_svg(input_svg_path, output_svg_path, watermark_text="DRAFT"):

    # Load existing SVG
    tree = ET.parse(input_svg_path)
    root = tree.getroot()

    # Ensure proper namespace
    if "xmlns" not in root.attrib:
        root.set("xmlns", "http://www.w3.org/2000/svg")
        root.set("xmlns:xlink", "http://www.w3.org/1999/xlink")

    # Get width and height of SVG (fallback if not set)
    width = int(root.get("width", 500))
    height = int(root.get("height", 500))

    # Add explicit white background (fix grey block issue in some viewers)
    background = ET.Element("rect", {
        "x": "0", "y": "0",
        "width": str(width),
        "height": str(height),
        "fill": "white"
    })
    root.insert(0, background)

    # Calculate center
    cx, cy = width // 2, height // 2

    # Create watermark text element with cross-compatible style
    watermark = ET.Element("text", {
        "x": str(cx),
        "y": str(cy),
        "font-size": "100",
        "fill": "black",
        "fill-opacity": "0.15",   # semi-transparent
        "text-anchor": "middle",  # center horizontally
        "transform": f"rotate(-30,{cx},{cy})",  # rotate around center
        # Use style instead of dominant-baseline for wider support
        "style": "alignment-baseline: middle; text-anchor: middle;"
    })
    watermark.text = watermark_text

    # Append watermark
    root.append(watermark)

    # Save new SVG with XML declaration
    tree.write(output_svg_path, encoding="utf-8", xml_declaration=True)

@celery_app.task
def add_watermark_to_pptx(input_file, output_file, watermark_text="DRAFT"):
    ext = os.path.splitext(input_file)[1].lower()

    if ext == ".ppt":
        print("Converting .ppt to .pptx...")
        input_file = convert_ppt_to_pptx(input_file)

    if not input_file.endswith(".pptx"):
        raise ValueError("File format not supported. Use .ppt or .pptx")

    add_watermark(input_file, output_file, watermark_text)
    